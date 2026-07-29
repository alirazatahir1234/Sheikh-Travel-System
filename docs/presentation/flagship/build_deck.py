#!/usr/bin/env python3
"""Build SheikhGo flagship executive PPTX from content/slides.yaml."""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
CONTENT = ROOT / "content" / "slides.yaml"
ASSETS = ROOT / "assets"
OUTPUT = ROOT / "output"
PPTX_NAME = "SheikhGo-AI-Fleet-Operations-Platform.pptx"

# Brand
EMERALD = RGBColor(0x0B, 0x6B, 0x50)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF8, 0xFA, 0xFC)
VISION = RGBColor(0x1E, 0x3A, 0x5F)
AMBER = RGBColor(0xB4, 0x53, 0x09)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)
LIGHT_NAVY = RGBColor(0xE0, 0xE7, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def hex_rgb(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def set_run(run, *, size=14, bold=False, color=NAVY, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, *, size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_shape(slide, shape_type, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def shape_text(shape, text, *, size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    # Clear existing
    if p.runs:
        p.runs[0].text = text
        set_run(p.runs[0], size=size, bold=bold, color=color)
    else:
        run = p.add_run()
        run.text = text
        set_run(run, size=size, bold=bold, color=color)
    try:
        shape.text_frame.paragraphs[0].space_before = Pt(0)
    except Exception:
        pass


def add_bullets(slide, left, top, width, height, items, *, size=15, color=SLATE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, color=color)
    return box


def chrome(slide, meta, slide_data, index, total, *, dark=False):
    """Left accent + footer."""
    accent = VISION if slide_data.get("badge") == "Vision" else EMERALD
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.12), SLIDE_H, fill=accent)
    footer_color = RGBColor(0x94, 0xA3, 0xB8) if not dark else RGBColor(0x94, 0xA3, 0xB8)
    add_textbox(
        slide,
        Inches(0.45),
        Inches(7.1),
        Inches(9),
        Inches(0.3),
        f"{meta.get('product', 'SheikhGo')}  ·  {meta.get('company', '')}",
        size=10,
        color=footer_color,
    )
    add_textbox(
        slide,
        Inches(11.2),
        Inches(7.1),
        Inches(1.7),
        Inches(0.3),
        f"{index} / {total}",
        size=10,
        color=footer_color,
        align=PP_ALIGN.RIGHT,
    )


def badge_pill(slide, badge: str | None):
    if not badge:
        return
    colors = {
        "Available": (EMERALD, WHITE),
        "Vision": (VISION, WHITE),
        "Indicative": (AMBER, WHITE),
        "Partial": (TEAL, WHITE),
    }
    fill, fg = colors.get(badge, (SLATE, WHITE))
    label = {"Vision": "Vision / Next-Gen", "Available": "Available Today", "Indicative": "Indicative", "Partial": "Partial"}.get(
        badge, badge
    )
    shape = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(0.28), Inches(2.5), Inches(0.36), fill=fill)
    shape.adjustments[0] = 0.5
    shape_text(shape, label, size=11, bold=True, color=fg)


def title_block(slide, title: str, *, top=Inches(0.35)):
    add_textbox(slide, Inches(0.5), top, Inches(9.5), Inches(0.55), title, size=28, bold=True, color=NAVY)


def logo(slide, *, white=False, left=Inches(0.5), top=Inches(0.35), width=Inches(1.6)):
    path = ASSETS / ("sheikhgo-logo-white.png" if white else "sheikhgo-logo.png")
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width)


def notes_for(slide, text: str | None):
    if not text:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ─── Layouts ─────────────────────────────────────────────────────────────────


def layout_cover(prs, meta, s, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H, fill=EMERALD)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(6.6), SLIDE_W, Inches(0.9), fill=EMERALD)
    logo(slide, white=True, left=Inches(0.7), top=Inches(0.55), width=Inches(2.0))
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(11.5), Inches(1.2), s["title"], size=36, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.7), Inches(3.6), Inches(11), Inches(0.5), s.get("subtitle", meta.get("subtitle", "")), size=20, color=LIGHT_TEAL)
    if s.get("bullets"):
        add_textbox(slide, Inches(0.7), Inches(4.4), Inches(11), Inches(0.4), s["bullets"][0], size=14, color=RGBColor(0xCB, 0xD5, 0xE1))
    add_textbox(slide, Inches(0.7), Inches(6.8), Inches(8), Inches(0.35), meta.get("company", ""), size=14, bold=True, color=WHITE)
    add_textbox(slide, Inches(10), Inches(6.8), Inches(2.8), Inches(0.35), meta.get("year", ""), size=14, color=WHITE, align=PP_ALIGN.RIGHT)
    notes_for(slide, s.get("notes"))


def layout_closing(prs, meta, s, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H, fill=EMERALD)
    logo(slide, white=True, left=Inches(0.7), top=Inches(0.55), width=Inches(2.0))
    add_textbox(slide, Inches(0.7), Inches(2.5), Inches(11.5), Inches(0.8), s["title"], size=40, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.7), Inches(3.4), Inches(11), Inches(0.45), s.get("subtitle", ""), size=18, color=LIGHT_TEAL)
    y = Inches(4.2)
    for b in s.get("bullets") or []:
        add_textbox(slide, Inches(0.7), y, Inches(11), Inches(0.35), b, size=15, color=RGBColor(0xCB, 0xD5, 0xE1))
        y = Emu(y + Inches(0.4))
    chrome(slide, meta, s, index, total, dark=True)
    notes_for(slide, s.get("notes"))


def layout_bullets(prs, meta, s, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill=SOFT)
    chrome(slide, meta, s, index, total)
    badge_pill(slide, s.get("badge"))
    title_block(slide, s["title"])
    top = Inches(1.15)
    if s.get("kpis"):
        draw_kpi_row(slide, s["kpis"], top=Inches(1.1))
        top = Inches(2.55)
    add_bullets(slide, Inches(0.55), top, Inches(12.2), Inches(4.2), s.get("bullets"), size=16)
    if s.get("footnote"):
        add_textbox(slide, Inches(0.55), Inches(6.7), Inches(12), Inches(0.3), s["footnote"], size=9, color=MUTED)
    notes_for(slide, s.get("notes"))


def layout_two_column(prs, meta, s, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill=SOFT)
    chrome(slide, meta, s, index, total)
    badge_pill(slide, s.get("badge"))
    title_block(slide, s["title"])

    left_box = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.9), Inches(5.4), fill=WHITE, line=LINE)
    right_box = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.2), Inches(5.9), Inches(5.4), fill=WHITE, line=LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.9), Inches(0.12), fill=EMERALD)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.85), Inches(1.2), Inches(5.9), Inches(0.12), fill=TEAL)

    add_textbox(slide, Inches(0.75), Inches(1.5), Inches(5.4), Inches(0.4), s.get("left_title", "Left"), size=16, bold=True, color=EMERALD)
    add_textbox(slide, Inches(7.1), Inches(1.5), Inches(5.4), Inches(0.4), s.get("right_title", "Right"), size=16, bold=True, color=TEAL)
    add_bullets(slide, Inches(0.75), Inches(2.05), Inches(5.4), Inches(4.2), s.get("left"), size=14)
    add_bullets(slide, Inches(7.1), Inches(2.05), Inches(5.4), Inches(4.2), s.get("right"), size=14)
    if s.get("footnote"):
        add_textbox(slide, Inches(0.55), Inches(6.7), Inches(12), Inches(0.3), s["footnote"], size=9, color=MUTED)
    notes_for(slide, s.get("notes"))
    _ = (left_box, right_box)


def layout_comparison(prs, meta, s, index, total):
    layout_two_column(prs, meta, s, index, total)


def draw_kpi_row(slide, kpis, *, top=Inches(1.15)):
    n = min(len(kpis), 4)
    if n == 0:
        return
    gap = Inches(0.25)
    total_w = Inches(12.3)
    card_w = Emu((total_w - gap * (n - 1)) / n)
    left = Inches(0.5)
    for kpi in kpis[:n]:
        card = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, card_w, Inches(1.2), fill=WHITE, line=LINE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, card_w, Inches(0.08), fill=EMERALD)
        add_textbox(slide, Emu(left + Inches(0.2)), Emu(top + Inches(0.25)), Emu(card_w - Inches(0.3)), Inches(0.45), str(kpi.get("value", "")), size=20, bold=True, color=NAVY)
        add_textbox(slide, Emu(left + Inches(0.2)), Emu(top + Inches(0.7)), Emu(card_w - Inches(0.3)), Inches(0.35), str(kpi.get("label", "")), size=11, color=MUTED)
        if kpi.get("hint"):
            add_textbox(slide, Emu(left + Inches(0.2)), Emu(top + Inches(0.95)), Emu(card_w - Inches(0.3)), Inches(0.2), str(kpi["hint"]), size=9, color=MUTED)
        left = Emu(left + card_w + gap)
        _ = card


def layout_kpi_cards(prs, meta, s, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill=SOFT)
    chrome(slide, meta, s, index, total)
    badge_pill(slide, s.get("badge"))
    title_block(slide, s["title"])
    draw_kpi_row(slide, s.get("kpis") or [], top=Inches(1.2))
    if s.get("bullets"):
        add_bullets(slide, Inches(0.55), Inches(2.8), Inches(12.2), Inches(3.5), s["bullets"], size=15)
    if s.get("footnote"):
        add_textbox(slide, Inches(0.55), Inches(6.7), Inches(12), Inches(0.3), s["footnote"], size=9, color=MUTED)
    notes_for(slide, s.get("notes"))


def layout_roadmap(prs, meta, s, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill=SOFT)
    chrome(slide, meta, s, index, total)
    badge_pill(slide, s.get("badge"))
    title_block(slide, s["title"])
    phases = s.get("phases") or []
    n = len(phases) or 1
    gap = Inches(0.3)
    card_w = Emu((Inches(12.3) - gap * (n - 1)) / n)
    left = Inches(0.5)
    colors = [EMERALD, TEAL, VISION]
    for i, phase in enumerate(phases):
        fill = colors[i % len(colors)]
        card = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.3), card_w, Inches(5.3), fill=WHITE, line=LINE)
        header = add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, Inches(1.3), card_w, Inches(0.7), fill=fill)
        shape_text(header, str(phase.get("name", "")), size=18, bold=True, color=WHITE)
        items = phase.get("items") or []
        add_bullets(slide, Emu(left + Inches(0.25)), Inches(2.3), Emu(card_w - Inches(0.4)), Inches(4.0), items, size=14)
        left = Emu(left + card_w + gap)
        _ = card
    notes_for(slide, s.get("notes"))


def box_node(slide, left, top, w, h, text, fill=EMERALD, size=11):
    shape = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, w, h, fill=fill)
    shape_text(shape, text, size=size, bold=True, color=WHITE)
    return shape


def layout_diagram(prs, meta, s, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill=SOFT)
    chrome(slide, meta, s, index, total)
    badge_pill(slide, s.get("badge"))
    title_block(slide, s["title"])
    kind = s.get("diagram", "platform_overview")
    drawers.get(kind, draw_platform_overview)(slide)
    if s.get("bullets"):
        add_bullets(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.4), s["bullets"], size=12)
    notes_for(slide, s.get("notes"))


def draw_platform_overview(slide):
    box_node(slide, Inches(0.6), Inches(1.4), Inches(3.5), Inches(1.1), "ERP Web\nAngular Command Center", EMERALD, 13)
    box_node(slide, Inches(4.9), Inches(1.4), Inches(3.5), Inches(1.1), "HTTPS API\nASP.NET Core · MediatR", TEAL, 13)
    box_node(slide, Inches(9.2), Inches(1.4), Inches(3.5), Inches(1.1), "Mobile Apps\nFlutter Driver · Fleet", VISION, 13)
    box_node(slide, Inches(1.2), Inches(3.3), Inches(2.4), Inches(0.9), "SQL Server", NAVY, 12)
    box_node(slide, Inches(3.9), Inches(3.3), Inches(2.4), Inches(0.9), "Redis*", SLATE, 12)
    box_node(slide, Inches(6.6), Inches(3.3), Inches(2.4), Inches(0.9), "Traccar GPS", EMERALD, 12)
    box_node(slide, Inches(9.3), Inches(3.3), Inches(2.4), Inches(0.9), "SignalR / SMTP", TEAL, 12)
    add_textbox(slide, Inches(0.6), Inches(4.4), Inches(12), Inches(0.3), "* Redis optional — fail-open cache", size=10, color=MUTED)


def draw_system_architecture(slide):
    draw_platform_overview(slide)


def draw_clean_architecture(slide):
    layers = [
        ("Presentation — ERP / Mobile / API Controllers", EMERALD),
        ("Application — MediatR Handlers · Use Cases", TEAL),
        ("Domain — Fleet · GPS · Notifications · Tenancy", VISION),
        ("Infrastructure — SQL · Dapper · Traccar · Redis · SMTP", NAVY),
    ]
    top = Inches(1.4)
    for label, color in layers:
        box_node(slide, Inches(1.5), top, Inches(10.3), Inches(0.85), label, color, 14)
        top = Emu(top + Inches(1.0))


def draw_azure(slide):
    box_node(slide, Inches(0.7), Inches(1.5), Inches(3.8), Inches(3.5), "Azure Edge\n\nApp Service / Containers\nAPI Gateway\nKey Vault", EMERALD, 13)
    box_node(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(3.5), "Data Plane\n\nAzure SQL\nRedis Cache\nBlob Storage", TEAL, 13)
    box_node(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(3.5), "Integrations\n\nTraccar\nSMTP / Push\nMonitoring", VISION, 13)


def draw_multi_tenant(slide):
    box_node(slide, Inches(4.5), Inches(1.35), Inches(4.3), Inches(0.8), "Platform Admin", NAVY, 14)
    for i, name in enumerate(["Tenant A", "Tenant B", "Tenant C"]):
        left = Inches(0.8 + i * 4.1)
        box_node(slide, left, Inches(2.5), Inches(3.7), Inches(2.5), f"{name}\n\nUsers · Roles · Fleet\nGPS · Notifications\nTenantId isolation", EMERALD if i == 0 else TEAL if i == 1 else VISION, 12)


def draw_notification_arch(slide):
    steps = ["Events", "Compose", "Preferences", "Dispatch", "Channels", "Inbox"]
    left = Inches(0.5)
    for i, step in enumerate(steps):
        box_node(slide, left, Inches(2.2), Inches(1.8), Inches(1.2), step, EMERALD if i % 2 == 0 else TEAL, 12)
        if i < len(steps) - 1:
            add_textbox(slide, Emu(left + Inches(1.75)), Inches(2.55), Inches(0.35), Inches(0.4), "→", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        left = Emu(left + Inches(2.1))
    box_node(slide, Inches(2.5), Inches(3.9), Inches(8.3), Inches(1.0), "SignalR realtime  ·  Email SMTP  ·  Retention / Archive jobs", VISION, 13)


def draw_mobile_arch(slide):
    box_node(slide, Inches(0.7), Inches(1.5), Inches(3.8), Inches(3.5), "Flutter Apps\n\nDriver · Fleet\nRBAC shells\nOffline outbox", EMERALD, 13)
    box_node(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(3.5), "Shared API\n\nJWT refresh\nPermissions\nSignalR / FCM", TEAL, 13)
    box_node(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(3.5), "Device Layer\n\nGPS · Biometrics\nHive · EN/AR\nStore builds", VISION, 13)


def draw_ai_notification_flow(slide):
    steps = [
        ("Ingest", "GPS · Maint · Trips"),
        ("Score", "Severity · Context"),
        ("Filter", "Suppress · Dedupe"),
        ("Route", "Role · Channel"),
        ("Act", "Escalate · Inbox"),
    ]
    left = Inches(0.45)
    for i, (title, sub) in enumerate(steps):
        box_node(slide, left, Inches(1.8), Inches(2.3), Inches(2.2), f"{title}\n\n{sub}", VISION if i == 4 else EMERALD if i % 2 == 0 else TEAL, 12)
        left = Emu(left + Inches(2.55))


def draw_ai_decision_engine(slide):
    box_node(slide, Inches(0.6), Inches(1.6), Inches(3.5), Inches(3.2), "Signals\n\nFleet · Driver\nGPS · Maint\nNotifications", EMERALD, 13)
    box_node(slide, Inches(4.9), Inches(1.6), Inches(3.5), Inches(3.2), "Decision Engine\n\nCorrelate\nRecommend\nExplain", VISION, 13)
    box_node(slide, Inches(9.2), Inches(1.6), Inches(3.5), Inches(3.2), "Actions\n\nDispatch\nWO · Coach\nHuman approve", TEAL, 13)


def draw_predictive_maintenance(slide):
    steps = ["History", "Telemetry", "Risk Model", "Work Order", "Parts / Shop"]
    left = Inches(0.5)
    for i, step in enumerate(steps):
        box_node(slide, left, Inches(2.0), Inches(2.2), Inches(1.5), step, EMERALD if i < 2 else VISION if i == 2 else TEAL, 12)
        left = Emu(left + Inches(2.5))


def draw_fleet_health(slide):
    box_node(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(1.6), "Fleet Health\n82 / 100", EMERALD, 18)
    factors = [("Compliance", EMERALD), ("Maintenance", TEAL), ("GPS Health", VISION), ("Utilization", NAVY)]
    left = Inches(0.8)
    for name, color in factors:
        box_node(slide, left, Inches(3.5), Inches(2.8), Inches(1.2), name, color, 13)
        left = Emu(left + Inches(3.1))


def draw_ai_learning(slide):
    box_node(slide, Inches(0.7), Inches(2.0), Inches(3.5), Inches(2.2), "Recommendations", EMERALD, 14)
    box_node(slide, Inches(4.9), Inches(2.0), Inches(3.5), Inches(2.2), "Accept / Reject\nFeedback", TEAL, 14)
    box_node(slide, Inches(9.2), Inches(2.0), Inches(3.5), Inches(2.2), "Model Improve\nPer Tenant", VISION, 14)


def draw_ai_ecosystem(slide):
    box_node(slide, Inches(4.5), Inches(1.4), Inches(4.3), Inches(0.9), "AI Control Plane", VISION, 14)
    box_node(slide, Inches(0.6), Inches(2.7), Inches(3.8), Inches(2.0), "Event Fabric\nGPS · Ops · Mobile", EMERALD, 13)
    box_node(slide, Inches(4.75), Inches(2.7), Inches(3.8), Inches(2.0), "Intelligence\nDecision · Notify · Learn", TEAL, 13)
    box_node(slide, Inches(8.9), Inches(2.7), Inches(3.8), Inches(2.0), "Experiences\nCopilot · Dashboards", NAVY, 13)


drawers = {
    "platform_overview": draw_platform_overview,
    "system_architecture": draw_system_architecture,
    "clean_architecture": draw_clean_architecture,
    "azure": draw_azure,
    "multi_tenant": draw_multi_tenant,
    "notification_arch": draw_notification_arch,
    "mobile_arch": draw_mobile_arch,
    "ai_notification_flow": draw_ai_notification_flow,
    "ai_decision_engine": draw_ai_decision_engine,
    "predictive_maintenance": draw_predictive_maintenance,
    "fleet_health": draw_fleet_health,
    "ai_learning": draw_ai_learning,
    "ai_ecosystem": draw_ai_ecosystem,
}


def layout_mockup(prs, meta, s, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill=SOFT)
    chrome(slide, meta, s, index, total)
    badge_pill(slide, s.get("badge"))
    title_block(slide, s["title"])

    # Stylized UI frame
    frame = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.15), Inches(7.4), Inches(5.4), fill=WHITE, line=LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(7.4), Inches(0.55), fill=NAVY)
    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(5), Inches(0.35), f"SheikhGo · {s.get('mockup', 'dashboard').replace('_', ' ').title()}", size=12, bold=True, color=WHITE)

    kind = s.get("mockup", "kpi")
    if kind == "notifications":
        rows = [("Critical", "Speeding · Vehicle 42", EMERALD), ("Warning", "Geofence exit · Depot B", AMBER), ("Info", "WO #118 scheduled", TEAL), ("System", "Device offline recovered", SLATE)]
        y = Inches(1.95)
        for title, body, color in rows:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), y, Inches(6.9), Inches(0.85), fill=SOFT, line=LINE)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), y, Inches(0.12), Inches(0.85), fill=color)
            add_textbox(slide, Inches(1.1), Emu(y + Inches(0.12)), Inches(6), Inches(0.3), title, size=12, bold=True, color=NAVY)
            add_textbox(slide, Inches(1.1), Emu(y + Inches(0.42)), Inches(6), Inches(0.3), body, size=11, color=MUTED)
            y = Emu(y + Inches(1.0))
    elif kind == "copilot":
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.0), Inches(6.7), Inches(1.1), fill=SOFT, line=LINE)
        add_textbox(slide, Inches(1.05), Inches(2.15), Inches(6.3), Inches(0.8), "Which vehicles need maintenance this week?", size=13, color=SLATE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(3.4), Inches(6.7), Inches(2.2), fill=LIGHT_TEAL, line=LINE)
        add_textbox(
            slide,
            Inches(1.05),
            Inches(3.55),
            Inches(6.3),
            Inches(1.9),
            "Copilot: 7 vehicles show elevated risk — 3 overdue service, 2 high idle + fault history, 2 license/compliance flags. Open Work Orders?",
            size=13,
            color=NAVY,
        )
    elif kind == "ai_center":
        tiles = [("Providers", "Configured"), ("Capabilities", "8 toggles"), ("Guardrails", "On"), ("Acceptance", "64%")]
        positions = [(0.75, 1.95), (4.2, 1.95), (0.75, 3.7), (4.2, 3.7)]
        for (x, y), (t, v) in zip(positions, tiles):
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.2), Inches(1.4), fill=SOFT, line=LINE)
            add_textbox(slide, Inches(x + 0.2), Inches(y + 0.3), Inches(2.8), Inches(0.35), t, size=12, color=MUTED)
            add_textbox(slide, Inches(x + 0.2), Inches(y + 0.7), Inches(2.8), Inches(0.4), v, size=18, bold=True, color=NAVY)
    else:
        # analytics / kpi default
        metrics = [("Active", "128"), ("Trips", "64"), ("Alerts", "12"), ("WO Open", "9")]
        x = Inches(0.75)
        for label, val in metrics:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(1.95), Inches(1.6), Inches(1.1), fill=SOFT, line=LINE)
            add_textbox(slide, Emu(x + Inches(0.1)), Inches(2.1), Inches(1.4), Inches(0.35), val, size=18, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)
            add_textbox(slide, Emu(x + Inches(0.1)), Inches(2.55), Inches(1.4), Inches(0.3), label, size=11, color=MUTED, align=PP_ALIGN.CENTER)
            x = Emu(x + Inches(1.75))
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.4), Inches(6.9), Inches(2.7), fill=SOFT, line=LINE)
        add_textbox(slide, Inches(1.0), Inches(4.4), Inches(6.4), Inches(0.5), "Stylized analytics canvas — live charts in product", size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # Side bullets
    side = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.15), Inches(4.6), Inches(5.4), fill=WHITE, line=LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.2), Inches(1.15), Inches(4.6), Inches(0.12), fill=EMERALD)
    add_textbox(slide, Inches(8.45), Inches(1.45), Inches(4.1), Inches(0.35), "Highlights", size=14, bold=True, color=EMERALD)
    add_bullets(slide, Inches(8.45), Inches(1.95), Inches(4.1), Inches(4.3), s.get("bullets"), size=13)
    _ = (frame, side)
    notes_for(slide, s.get("notes"))


LAYOUTS = {
    "cover": layout_cover,
    "closing": layout_closing,
    "bullets": layout_bullets,
    "two_column": layout_two_column,
    "comparison": layout_comparison,
    "kpi_cards": layout_kpi_cards,
    "roadmap": layout_roadmap,
    "diagram": layout_diagram,
    "mockup": layout_mockup,
    "section": layout_bullets,
}


def try_export_pdf(pptx_path: Path) -> Path | None:
    pdf_path = OUTPUT / pptx_path.with_suffix(".pdf").name
    # LibreOffice
    for bin_name in ("soffice", "libreoffice"):
        try:
            subprocess.run(
                [bin_name, "--headless", "--convert-to", "pdf", "--outdir", str(OUTPUT), str(pptx_path)],
                check=True,
                capture_output=True,
                timeout=180,
            )
            if pdf_path.exists():
                return pdf_path
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
            pass

    # macOS PowerPoint via AppleScript
    script = ROOT / "export_pdf.applescript"
    if script.exists() and Path("/Applications/Microsoft PowerPoint.app").exists():
        try:
            subprocess.run(["osascript", str(script)], check=True, capture_output=True, timeout=300)
            if pdf_path.exists():
                return pdf_path
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
            pass

    # Fallback: HTML companion printable deck for PDF via browser print
    return None


def write_html_fallback(meta: dict, slides: list[dict], pptx_path: Path) -> Path:
    """Browser-printable companion when LibreOffice/PowerPoint unavailable."""
    html_path = OUTPUT / "SheikhGo-AI-Fleet-Operations-Platform.html"
    parts = [
        "<!DOCTYPE html><html><head><meta charset='utf-8'/>",
        "<title>{}</title>".format(meta.get("title", "SheikhGo")),
        "<style>",
        "@page { size: landscape; margin: 12mm; }",
        "body{font-family:Calibri,Segoe UI,sans-serif;background:#0f172a;color:#0f172a;margin:0;}",
        ".slide{background:#f8fafc;width:1100px;min-height:620px;margin:24px auto;padding:40px 48px;page-break-after:always;position:relative;border-left:8px solid #0B6B50;}",
        ".slide.vision{border-left-color:#1e3a5f;}",
        ".badge{display:inline-block;background:#0B6B50;color:#fff;padding:4px 10px;border-radius:999px;font-size:12px;font-weight:700;}",
        ".badge.Vision{background:#1e3a5f;}",
        ".badge.Indicative{background:#b45309;}",
        "h1{font-size:28px;margin:12px 0 16px;}",
        "ul{line-height:1.55;color:#475569;}",
        ".kpis{display:grid;grid-template-columns:repeat(4,1fr);gap:12px;margin:16px 0;}",
        ".kpi{background:#fff;border:1px solid #e2e8f0;padding:14px;border-top:4px solid #0B6B50;}",
        ".kpi b{display:block;font-size:20px;color:#0f172a;}",
        ".cols{display:grid;grid-template-columns:1fr 1fr;gap:16px;}",
        ".card{background:#fff;border:1px solid #e2e8f0;padding:16px;}",
        ".cover{background:#0f172a;color:#fff;border-left-color:#0B6B50;}",
        ".cover h1{color:#fff;font-size:36px;}",
        ".foot{position:absolute;bottom:16px;left:48px;right:48px;display:flex;justify-content:space-between;color:#94a3b8;font-size:12px;}",
        "@media print{body{background:#fff;} .slide{margin:0;width:auto;box-shadow:none;}}",
        "</style></head><body>",
    ]
    total = len(slides)
    for i, s in enumerate(slides, 1):
        badge = s.get("badge") or ""
        cls = "slide cover" if s.get("layout") in ("cover", "closing") else ("slide vision" if badge == "Vision" else "slide")
        parts.append(f"<div class='{cls}'>")
        if badge and s.get("layout") not in ("cover", "closing"):
            parts.append(f"<span class='badge {badge}'>{badge}</span>")
        parts.append(f"<h1>{s.get('title','')}</h1>")
        if s.get("subtitle"):
            parts.append(f"<p><em>{s['subtitle']}</em></p>")
        if s.get("kpis"):
            parts.append("<div class='kpis'>")
            for k in s["kpis"]:
                parts.append(f"<div class='kpi'><b>{k.get('value','')}</b>{k.get('label','')}<div style='font-size:11px;color:#94a3b8'>{k.get('hint','')}</div></div>")
            parts.append("</div>")
        if s.get("layout") in ("two_column", "comparison"):
            parts.append("<div class='cols'>")
            parts.append(f"<div class='card'><h3>{s.get('left_title','')}</h3><ul>")
            for b in s.get("left") or []:
                parts.append(f"<li>{b}</li>")
            parts.append("</ul></div>")
            parts.append(f"<div class='card'><h3>{s.get('right_title','')}</h3><ul>")
            for b in s.get("right") or []:
                parts.append(f"<li>{b}</li>")
            parts.append("</ul></div></div>")
        if s.get("phases"):
            parts.append("<div class='cols'>")
            for ph in s["phases"]:
                parts.append(f"<div class='card'><h3>{ph.get('name')}</h3><ul>")
                for item in ph.get("items") or []:
                    parts.append(f"<li>{item}</li>")
                parts.append("</ul></div>")
            parts.append("</div>")
        if s.get("bullets"):
            parts.append("<ul>")
            for b in s["bullets"]:
                parts.append(f"<li>{b}</li>")
            parts.append("</ul>")
        if s.get("footnote"):
            parts.append(f"<p style='font-size:11px;color:#94a3b8'>{s['footnote']}</p>")
        parts.append(f"<div class='foot'><span>{meta.get('product')} · {meta.get('company')}</span><span>{i} / {total}</span></div>")
        parts.append("</div>")
    parts.append("</body></html>")
    html_path.write_text("\n".join(parts), encoding="utf-8")
    return html_path


def build():
    data = yaml.safe_load(CONTENT.read_text(encoding="utf-8"))
    meta = data.get("meta") or {}
    slides = data.get("slides") or []
    if len(slides) != 60:
        print(f"WARNING: expected 60 slides, found {len(slides)}", file=sys.stderr)

    OUTPUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(slides)
    for i, s in enumerate(slides, 1):
        layout = s.get("layout", "bullets")
        fn = LAYOUTS.get(layout, layout_bullets)
        fn(prs, meta, s, i, total)

    out = OUTPUT / PPTX_NAME
    prs.save(str(out))
    print(f"Wrote {out} ({total} slides)")

    pdf = try_export_pdf(out)
    if pdf:
        print(f"Wrote {pdf}")
    else:
        html = write_html_fallback(meta, slides, out)
        print(f"PDF converter not found — wrote printable HTML companion: {html}")
        print("Open in browser and Print → Save as PDF, or install LibreOffice / use PowerPoint Export.")
        # Also try macOS cupsfilter / textutil won't help for pptx.
        # Create a simple note file
        (OUTPUT / "PDF-EXPORT.txt").write_text(
            "LibreOffice/PowerPoint not detected.\n"
            "Options:\n"
            "1) Open the .pptx in PowerPoint/Keynote → Export → PDF\n"
            "2) Open SheikhGo-AI-Fleet-Operations-Platform.html → Print → Save as PDF\n"
            "3) soffice --headless --convert-to pdf --outdir output output/*.pptx\n",
            encoding="utf-8",
        )


if __name__ == "__main__":
    build()
